import os
import win32com.client
from pptx import Presentation


def convert_ppt_to_pptx(file_path):
    # Create a PowerPoint application instance
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")

    # Open the .ppt file
    presentation = powerpoint.Presentations.Open(os.path.abspath(file_path))

    # Initialize an empty string to store the extracted text
    extracted_text = ""

    # Loop through slides and extract text
    for slide in presentation.Slides:
        for shape in slide.Shapes:
            if shape.HasTextFrame:
                extracted_text += shape.TextFrame.TextRange.Text + '\n'

    # Close the presentation
    presentation.Close()

    # Quit PowerPoint application
    powerpoint.Quit()

    return extracted_text


def ppt2txt(file_path, file_extension):
    extracted_text = ''
    if file_extension == 'ppt':
        extracted_text = convert_ppt_to_pptx(file_path)
    else:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    extracted_text += shape.text + '\n'
    return extracted_text
